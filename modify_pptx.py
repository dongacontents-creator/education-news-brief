#!/usr/bin/env python3
"""Modify the education news brief PPTX presentation."""

import copy
from pptx import Presentation
from pptx.util import Pt, Emu, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn, nsmap
from pptx.oxml import parse_xml
from lxml import etree

SRC = '/root/.claude/uploads/a7b122b2-ff33-44f0-8e4f-28a38616b184/8f262ace-_________________________.pptx'
DST = '/home/user/education-news-brief/교육뉴스브리핑_개발현황_v2.pptx'

# Colors
RED = RGBColor(0xC1, 0x12, 0x1F)
RED_ACCENT = RGBColor(0xC0, 0x00, 0x00)
BADGE_RED = RGBColor(0xC0, 0x39, 0x2B)
DARK = RGBColor(0x1E, 0x29, 0x3B)
BODY = RGBColor(0x47, 0x55, 0x69)
SUBTITLE = RGBColor(0x64, 0x74, 0x8B)
PILL_BG = RGBColor(0xFE, 0xE2, 0xE2)
PILL_TEXT = RGBColor(0x99, 0x1B, 0x1B)
SOL_BG = RGBColor(0xFF, 0xF1, 0xF2)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
SHADOW = RGBColor(0xC8, 0xCD, 0xD4)

W = 12192000
H = 6858000


def add_solid_shape(slide, left, top, width, height, color, shape_type=1):
    """Add a rectangle with solid fill."""
    from pptx.util import Emu
    shape = slide.shapes.add_shape(shape_type, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_textbox(slide, left, top, width, height):
    from pptx.util import Emu
    txBox = slide.shapes.add_textbox(left, top, width, height)
    return txBox


def set_run(run, text, color, size_emu, bold=False):
    run.text = text
    run.font.color.rgb = color
    run.font.size = size_emu
    run.font.bold = bold
    run.font.name = '맑은 고딕'


def add_title_textbox(slide, section_num, title_text, subtitle_text=''):
    """Add the standard title bar for a slide."""
    # Title
    tb = slide.shapes.add_textbox(571499, 476402, 10000000, 467258)
    tf = tb.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    r1 = p.add_run()
    set_run(r1, f'{section_num}. ', RED_ACCENT, 317500, bold=True)
    r2 = p.add_run()
    set_run(r2, title_text, DARK, 317500, bold=True)

    # Subtitle
    if subtitle_text:
        tb2 = slide.shapes.add_textbox(571500, 984105, 11000000, 277063)
        tf2 = tb2.text_frame
        p2 = tf2.paragraphs[0]
        r = p2.add_run()
        set_run(r, subtitle_text, SUBTITLE, 190500)
    return tb


def add_top_bar(slide):
    """Add the red top bar."""
    bar = slide.shapes.add_shape(1, 0, 0, W, 152705)
    bar.fill.solid()
    bar.fill.fore_color.rgb = RED
    bar.line.fill.background()
    return bar


def add_card(slide, left, top, width, height, border_color=None):
    """Add a white card with optional left border."""
    card = slide.shapes.add_shape(1, left, top, width, height)
    card.fill.solid()
    card.fill.fore_color.rgb = WHITE
    card.line.fill.background()

    if border_color:
        border_w = 38405
        border = slide.shapes.add_shape(1, left, top, border_w, height)
        border.fill.solid()
        border.fill.fore_color.rgb = border_color
        border.line.fill.background()
    return card


def add_solution_box(slide, left, top, width, height):
    """Add a solution box with light pink bg."""
    box = slide.shapes.add_shape(1, left, top, width, height)
    box.fill.solid()
    box.fill.fore_color.rgb = SOL_BG
    box.line.fill.background()
    return box


def add_step_badge(slide, left, top, num_text):
    """Add step number badge (dark red)."""
    badge = slide.shapes.add_shape(1, left, top, 351525, 234350)
    badge.fill.solid()
    badge.fill.fore_color.rgb = BADGE_RED
    badge.line.fill.background()

    tb = slide.shapes.add_textbox(left, top, 351525, 234350)
    tf = tb.text_frame
    tf.auto_size = None
    from pptx.enum.text import PP_ALIGN
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    r = tf.paragraphs[0].add_run()
    set_run(r, num_text, WHITE, 152400, bold=True)
    # vertical center
    from pptx.enum.text import MSO_ANCHOR
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    return badge


def add_label_pill(slide, left, top, label_text):
    """Add keyword label pill."""
    pill_w = max(656180, len(label_text) * 120000 + 200000)
    pill = slide.shapes.add_shape(1, left, top, pill_w, 234350)
    pill.fill.solid()
    pill.fill.fore_color.rgb = PILL_BG
    pill.line.fill.background()

    tb = slide.shapes.add_textbox(left + 50000, top, pill_w - 100000, 234350)
    tf = tb.text_frame
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    r = tf.paragraphs[0].add_run()
    set_run(r, label_text, PILL_TEXT, 133350, bold=True)
    return pill


def add_body_text(slide, left, top, width, height, text, color=None, size=None, bold=False, align=None):
    """Add a text box with body text."""
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    if align:
        p.alignment = align
    r = p.add_run()
    c = color or BODY
    s = size or 152400
    set_run(r, text, c, s, bold=bold)
    return tb


def add_multiline_body(slide, left, top, width, height, lines, color=None, size=None):
    """Add a text box with multiple lines."""
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        r = p.add_run()
        c = color or BODY
        s = size or 152400
        set_run(r, line, c, s)
    return tb


def move_slide(prs, old_idx, new_idx):
    """Move a slide from old_idx to new_idx in the presentation."""
    xml_slides = prs.slides._sldIdLst
    slides = list(xml_slides)
    el = slides[old_idx]
    xml_slides.remove(el)
    if new_idx >= len(list(xml_slides)):
        xml_slides.append(el)
    else:
        xml_slides.insert(new_idx, el)


def add_blank_slide(prs):
    """Add a blank slide using the first layout."""
    layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(layout)
    # Remove all placeholder shapes from layout
    for ph in slide.placeholders:
        sp = ph._element
        sp.getparent().remove(sp)
    return slide


# ============================================================
# BUILD THE PRESENTATION
# ============================================================
prs = Presentation(SRC)
n_orig = len(prs.slides)
print(f'Original slide count: {n_orig}')

# ============================================================
# 1. MODIFY SLIDE 3 - Add timing and details
# ============================================================
slide3 = prs.slides[2]

# Add timing/background info box at the bottom of slide 3
# Adding a new card below the existing content with timing info
timing_card = add_card(slide3, 571499, 5900000, 11000000, 700000, border_color=RED)
add_body_text(
    slide3, 630000, 5950000, 10800000, 200000,
    '⏱ 개발 기간 및 방식', color=DARK, size=190500, bold=True
)
add_body_text(
    slide3, 630000, 6150000, 10800000, 400000,
    '2026년 3월 ~ 현재 (약 8주 진행 중)  |  바이브코딩(Vibe Coding) 방식 적용 — AI 코드 에디터(Claude Code)를 활용한 자연어 기반 개발',
    color=SUBTITLE, size=152400
)

print('Slide 3 updated with timing info')

# ============================================================
# 2. ADD 5 STEP SLIDES after slide 4
# ============================================================
steps = [
    {
        'num': '01',
        'label': 'INPUT',
        'heading': 'Claude에 자연어로 요구 사항 설명',
        'desc': [
            '기능·화면 구조 관련 요구 사항을 자연어로 상세히 작성',
            '',
            '• 구현하려는 기능을 일상 언어로 구체적으로 설명',
            '• 화면 레이아웃, 데이터 흐름, 동작 방식 등을 자세히 기술',
            '• 기존 코드와의 연계 지점, 제약 조건도 함께 명시',
        ],
        'note': '요구 사항이 구체적일수록 AI의 결과물 품질이 높아짐 → 자연어 설명 역량이 개발 효율을 결정하는 핵심 요소',
        'tool': 'Claude (claude.ai)',
        'duration': '10~30분',
    },
    {
        'num': '02',
        'label': 'GENERATE',
        'heading': 'Claude로 1차 결과물 도출',
        'desc': [
            '코드·구조·로직 초안을 Claude가 자동 생성',
            '',
            '• 요구 사항 기반으로 HTML/CSS/JavaScript/Python 코드 자동 생성',
            '• 전체 파일 구조, 함수 설계, 데이터 모델 초안 제시',
            '• 여러 구현 방안을 비교 검토하여 최적 방향 선택',
        ],
        'note': '초기(w1~w7)에는 Claude가 생성한 HTML을 수동으로 GitHub에 업로드하는 방식 → 반복 수정 사이클 소요',
        'tool': 'Claude (claude.ai)',
        'duration': '30분~2시간',
    },
    {
        'num': '03',
        'label': 'REFINE',
        'heading': 'Claude Code로 코드 수정 및 기능 구현',
        'desc': [
            '네이버 뉴스 API 수집 적용 / Groq AI 요약 및 분석 기능 연동',
            '',
            '• Claude Code(터미널 기반) 활용 → 실제 파일에 직접 코드 반영',
            '• 네이버 뉴스 API 연동: 키워드 검색, 날짜 필터, 출처 파싱 구현',
            '• Groq AI 연동: 기사별 한 줄 요약, 교육 분야 분류, 주간 인사이트 생성',
            '• GitHub Secrets로 API 키 보안 관리',
        ],
        'note': '현행(w8~) 핵심 단계 — 자동화 파이프라인의 실질적 구현 단계',
        'tool': 'Claude Code (CLI)',
        'duration': '1~4시간',
    },
    {
        'num': '04',
        'label': 'DEPLOY',
        'heading': 'GitHub Actions로 자동 수집·배포',
        'desc': [
            'Actions workflow 자동 실행(매주 일요일 23:00) / GitHub Pages 자동 배포',
            '',
            '• collect_news.py: 뉴스 수집 → weeks.json 업데이트',
            '• generate_site.py: weeks.json → index.html 자동 생성',
            '• GitHub Pages로 정적 사이트 자동 배포',
            '• 별도 배포 workflow(deploy_only.yml)로 긴급 수동 배포 지원',
        ],
        'note': 'GitHub Actions cron 스케줄로 매주 자동 실행 → 담당자 개입 없이 사이트 갱신',
        'tool': 'GitHub Actions / Pages',
        'duration': '5~15분 (자동)',
    },
    {
        'num': '05',
        'label': 'ITERATE',
        'heading': '오류 발생 시 ①~③ 반복',
        'desc': [
            '오류 발생 → 자연어 재설명 → 페이지 재생성의 반복 사이클',
            '',
            '• GitHub Actions 실패 로그를 Claude Code에 붙여넣기 → 오류 원인 분석',
            '• 수정 방향을 자연어로 설명 → 코드 재생성 → 재배포',
            '• 동일 오류 재발 방지를 위한 방어 로직(날짜 필터, 인코딩 변환 등) 추가',
            '• 필요 시 CLAUDE.md에 제약 사항 문서화하여 지속적 컨텍스트 유지',
        ],
        'note': '바이브코딩의 핵심 — 완벽한 코드보다 빠른 반복이 효율적. 오류도 자연어로 해결',
        'tool': 'Claude Code + GitHub Actions',
        'duration': '30분~2시간',
    },
]

for i, step in enumerate(steps):
    slide = add_blank_slide(prs)
    add_top_bar(slide)

    # Title
    tb = slide.shapes.add_textbox(571499, 476402, 11000000, 467258)
    tf = tb.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    r1 = p.add_run(); set_run(r1, '02.', RED_ACCENT, 317500, bold=True)
    r2 = p.add_run(); set_run(r2, ' 진행 과정 및 시행착오 - ① 전체 개발 과정 / ', DARK, 317500, bold=True)
    r3 = p.add_run(); set_run(r3, f'{step["num"]}. {step["heading"][:20]}', RED_ACCENT, 317500, bold=True)

    # Subtitle
    tb2 = slide.shapes.add_textbox(571500, 984105, 11000000, 277063)
    p2 = tb2.text_frame.paragraphs[0]
    r = p2.add_run(); set_run(r, f'Step {step["num"]} — {step["label"]} 단계 상세', SUBTITLE, 190500)

    # Step badge (top left area)
    add_step_badge(slide, 571499, 1400000, step['num'])

    # Label pill
    add_label_pill(slide, 950000, 1400000, step['label'])

    # Tool & Duration info
    add_body_text(slide, 571499 + 351525 + 50000 + 700000 + 100000, 1400000,
                  3000000, 240000,
                  f'🔧 {step["tool"]}   ⏱ {step["duration"]}',
                  color=SUBTITLE, size=127000)

    # Main card - description
    card_left = 571499
    card_top = 1800000
    card_w = 7200000
    card_h = 4700000
    add_card(slide, card_left, card_top, card_w, card_h, border_color=RED)

    # Card content
    inner_left = card_left + 38405 + 80000
    inner_top = card_top + 120000
    tb3 = slide.shapes.add_textbox(inner_left, inner_top, card_w - 38405 - 160000, card_h - 240000)
    tf3 = tb3.text_frame
    tf3.word_wrap = True
    for j, line in enumerate(step['desc']):
        if j == 0:
            p_line = tf3.paragraphs[0]
        else:
            p_line = tf3.add_paragraph()
        r_line = p_line.add_run()
        if j == 0:
            set_run(r_line, line, DARK, 171450, bold=True)
        else:
            set_run(r_line, line, BODY, 152400)

    # Note box (right side)
    note_left = card_left + card_w + 100000
    note_w = W - note_left - 200000
    note_top = card_top
    note_h = 1800000

    note_box = add_solution_box(slide, note_left, note_top, note_w, note_h)
    add_body_text(slide, note_left + 80000, note_top + 80000,
                  note_w - 160000, 200000,
                  '💡 인사이트', color=RED_ACCENT, size=171450, bold=True)
    add_multiline_body(slide, note_left + 80000, note_top + 300000,
                       note_w - 160000, note_h - 380000,
                       [step['note']], color=BODY, size=133350)

    # Bottom bar
    bottom_bar = slide.shapes.add_shape(1, 0, H - 100000, W, 100000)
    bottom_bar.fill.solid()
    bottom_bar.fill.fore_color.rgb = RED
    bottom_bar.line.fill.background()

    print(f'Step slide {i+1} ({step["label"]}) created at index {len(prs.slides)-1}')

# Move the 5 step slides to positions 5-9 (after original slide 4)
# Currently they are at positions n_orig through n_orig+4
# We want them at positions 4, 5, 6, 7, 8 (0-indexed)
for i in range(5):
    # After adding i slides, step slide i is at: n_orig + i - i = n_orig (still at n_orig each time after moving)
    # Actually each time we move n_orig to position 4+i
    move_slide(prs, n_orig + i, 4 + i)
    print(f'Moved step slide {i+1} to position {4+i}')

print(f'Current slide count: {len(prs.slides)}')

# ============================================================
# 3. ADD 3 ERROR SLIDES after original slide 5 (now at index 9)
# ============================================================
errors = [
    {
        'num': '1',
        'title': '문제1. 기사 수집 날짜 오류',
        'problem': '최근 1주간의 뉴스를 요청했음에도 작년 뉴스가 다수 섞이는 오류 발생\n또는 동일한 주차 데이터가 중복으로 생성되는 오류 발생',
        'cause': '네이버 뉴스 API의 날짜 파라미터가 완전히 신뢰할 수 없는 구조\n기사 발행 날짜와 API 인덱스 날짜 불일치로 인해 과거 기사 포함',
        'solution': [
            '해결 완료',
            '수집 주차의 월~일 날짜 범위 필터 추가',
            '해당 날짜 범위 외 기사 자동 제외',
            '주차 판별 로직: isoweek 기반 월요일~일요일 범위 계산',
        ],
    },
    {
        'num': '2',
        'title': '문제2. 기사 제목 인코딩 오류',
        'problem': '기사 제목에 포함된 따옴표(")가 &quot; 형태로 화면에 그대로 노출\n특수문자(< > & " \')가 HTML 인코딩 상태로 카드에 표시',
        'cause': '네이버 뉴스 API가 기사 제목을 HTML 인코딩된 형태로 전달\nJSON 저장 단계에서 디코딩 처리 없이 그대로 저장',
        'solution': [
            '해결 완료',
            '제목 수집 함수에 인코딩 자동 변환 추가',
            '수집 시 &quot; → " 자동 변환되도록 수정',
            'html.unescape() 함수로 모든 HTML 엔터티 자동 변환',
        ],
    },
    {
        'num': '3',
        'title': '문제3. 기사 자동 수집 기능 오류',
        'problem': '5월 1주차 성공 이후 신규 기능 추가를 위한 코드 수정 과정에서 5월 2주차 자동 수집 실패\n인사이트 기능 추가 중 NameError 발생된 상태로 GitHub에 반영',
        'cause': 'strip_tags 함수가 인사이트 모듈 추가 과정에서 누락\n오류 발생 코드가 테스트 없이 배포되어 전체 workflow 실패',
        'solution': [
            '해결 완료',
            '오류 발생 코드 확인 후 strip_tags 함수 복구',
            '자동 수집 관련 코드 재검토 및 안정화 시도',
            '5월 3주차 자동 수집 기능 정상화 완료',
            '이후 코드 수정 시 로컬 테스트 → GitHub 반영 프로세스 정착',
        ],
    },
]

for i, err in enumerate(errors):
    slide = add_blank_slide(prs)
    add_top_bar(slide)

    # Title
    tb = slide.shapes.add_textbox(571499, 476402, 11000000, 467258)
    tf = tb.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    r1 = p.add_run(); set_run(r1, '02.', RED_ACCENT, 317500, bold=True)
    r2 = p.add_run(); set_run(r2, ' 진행 과정 및 시행착오 - ② 주요 버그 발견 및 해결 과정 / ', DARK, 317500, bold=True)
    r3 = p.add_run(); set_run(r3, f'문제{err["num"]}', RED_ACCENT, 317500, bold=True)

    # Subtitle
    tb2 = slide.shapes.add_textbox(571500, 984105, 11000000, 277063)
    p2 = tb2.text_frame.paragraphs[0]
    r = p2.add_run(); set_run(r, err['title'], SUBTITLE, 190500)

    # Problem section header
    add_body_text(slide, 571499, 1400000, 2000000, 250000, '▶ 문제 상황', color=DARK, size=190500, bold=True)

    # Problem card
    add_card(slide, 571499, 1650000, 5500000, 1500000, border_color=RED)
    tb_prob = slide.shapes.add_textbox(571499 + 38405 + 80000, 1720000, 5200000, 1360000)
    tf_prob = tb_prob.text_frame
    tf_prob.word_wrap = True
    for j, line in enumerate(err['problem'].split('\n')):
        if j == 0:
            p_line = tf_prob.paragraphs[0]
        else:
            p_line = tf_prob.add_paragraph()
        r_line = p_line.add_run()
        set_run(r_line, line, BODY, 152400)

    # Cause section
    add_body_text(slide, 571499, 3250000, 2000000, 250000, '▶ 발생 원인', color=DARK, size=190500, bold=True)

    add_card(slide, 571499, 3500000, 5500000, 1200000, border_color=RGBColor(0x71, 0x80, 0x96))
    tb_cause = slide.shapes.add_textbox(571499 + 38405 + 80000, 3570000, 5200000, 1060000)
    tf_cause = tb_cause.text_frame
    tf_cause.word_wrap = True
    for j, line in enumerate(err['cause'].split('\n')):
        if j == 0:
            p_line = tf_cause.paragraphs[0]
        else:
            p_line = tf_cause.add_paragraph()
        r_line = p_line.add_run()
        set_run(r_line, line, SUBTITLE, 152400)

    # Solution box (right side)
    sol_left = 6300000
    sol_w = W - sol_left - 200000
    sol_top = 1400000
    sol_h = 5000000
    add_solution_box(slide, sol_left, sol_top, sol_w, sol_h)

    add_body_text(slide, sol_left + 80000, sol_top + 80000,
                  sol_w - 160000, 250000,
                  '✅ 해결 방법', color=RED_ACCENT, size=190500, bold=True)

    tb_sol = slide.shapes.add_textbox(sol_left + 80000, sol_top + 380000, sol_w - 160000, sol_h - 460000)
    tf_sol = tb_sol.text_frame
    tf_sol.word_wrap = True
    for j, line in enumerate(err['solution']):
        if j == 0:
            p_line = tf_sol.paragraphs[0]
        else:
            p_line = tf_sol.add_paragraph()
        r_line = p_line.add_run()
        if j == 0:
            set_run(r_line, '☑ ' + line, RED_ACCENT, 171450, bold=True)
        else:
            set_run(r_line, '• ' + line, BODY, 152400)

    # Bottom bar
    bottom_bar = slide.shapes.add_shape(1, 0, H - 100000, W, 100000)
    bottom_bar.fill.solid()
    bottom_bar.fill.fore_color.rgb = RED
    bottom_bar.line.fill.background()

    print(f'Error slide {i+1} (문제{err["num"]}) created at index {len(prs.slides)-1}')

# Now move the 3 error slides to positions after original slide 5
# Original slide 5 (버그) is now at index 9 (0-indexed: 0,1,2,3, 5step slides=4-8, orig5=9)
# We want error slides at 10, 11, 12
orig_slide5_idx = 9  # 0-indexed
n_before_errors = len(prs.slides) - 3  # index where first error slide currently is

for i in range(3):
    move_slide(prs, n_before_errors + i, orig_slide5_idx + 1 + i)
    print(f'Moved error slide {i+1} to position {orig_slide5_idx + 1 + i}')

print(f'Current slide count: {len(prs.slides)}')

# ============================================================
# 4. ADD INSIGHT → CURRENT WORK SLIDE
# ============================================================
insight_slide = add_blank_slide(prs)
add_top_bar(insight_slide)

# Title
tb = insight_slide.shapes.add_textbox(571499, 476402, 11000000, 467258)
tf = tb.text_frame
tf.word_wrap = False
p = tf.paragraphs[0]
r1 = p.add_run(); set_run(r1, '05.', RED_ACCENT, 317500, bold=True)
r2 = p.add_run(); set_run(r2, ' 도출된 인사이트의 현업 적용', DARK, 317500, bold=True)

# Subtitle
tb2 = insight_slide.shapes.add_textbox(571500, 984105, 11000000, 277063)
p2 = tb2.text_frame.paragraphs[0]
r = p2.add_run()
set_run(r, '프로젝트를 통해 얻은 인사이트를 현재 업무에 적용하고 있는 사례', SUBTITLE, 190500)

# 3 insight cards
insights = [
    {
        'icon': '🤖',
        'title': 'AI 협업 방식 표준화',
        'body': [
            '바이브코딩 경험을 통해 팀 내 AI 활용 가이드라인 정립',
            '• 요구 사항 작성 템플릿 내재화 (기능/화면/제약 3단계)',
            '• Claude Code 터미널 방식으로 직접 파일 수정 → 실수 최소화',
            '• CLAUDE.md 활용: 반복 실수 방지 제약 사항 문서화',
        ],
    },
    {
        'icon': '📊',
        'title': '데이터 기반 콘텐츠 기획',
        'body': [
            '주간 수집 데이터를 실제 콘텐츠 기획 회의에서 활용',
            '• 교육 정책 변화 트렌드 파악 → 교과서·학습 콘텐츠 방향 검토',
            '• AI 교육 동향 주간 리포트 → 서비스 기획팀 공유 자료로 활용',
            '• 키워드 빈도 분석 → 경쟁사 모니터링 인사이트 도출',
        ],
    },
    {
        'icon': '⚙️',
        'title': '자동화 사고방식 업무 적용',
        'body': [
            '반복 업무의 자동화 가능성을 우선 검토하는 방식으로 사고 전환',
            '• 매주 수동 수집 → 자동 파이프라인 구축으로 주당 3~4시간 절감',
            '• GitHub Actions 활용 패턴을 다른 반복 업무에도 적용 검토',
            '• "AI에게 설명 가능한 업무"부터 자동화 시도하는 원칙 확립',
        ],
    },
]

card_w = 3600000
card_h = 4400000
card_gap = 190000
for i, ins in enumerate(insights):
    left = 571499 + i * (card_w + card_gap)
    top = 1500000

    add_card(insight_slide, left, top, card_w, card_h, border_color=RED)

    # Icon + title
    add_body_text(insight_slide, left + 38405 + 80000, top + 120000,
                  card_w - 200000, 300000,
                  f'{ins["icon"]} {ins["title"]}', color=DARK, size=190500, bold=True)

    # Divider
    divider = insight_slide.shapes.add_shape(1, left + 38405 + 80000, top + 460000,
                                              card_w - 280000, 19050)
    divider.fill.solid()
    divider.fill.fore_color.rgb = RGBColor(0xE2, 0xE8, 0xF0)
    divider.line.fill.background()

    # Body text
    tb_ins = insight_slide.shapes.add_textbox(left + 38405 + 80000, top + 510000,
                                               card_w - 280000, card_h - 590000)
    tf_ins = tb_ins.text_frame
    tf_ins.word_wrap = True
    for j, line in enumerate(ins['body']):
        if j == 0:
            p_line = tf_ins.paragraphs[0]
        else:
            p_line = tf_ins.add_paragraph()
        r_line = p_line.add_run()
        if j == 0:
            set_run(r_line, line, SUBTITLE, 152400, bold=True)
        else:
            set_run(r_line, line, BODY, 133350)

# Summary bar at bottom
sum_bar = insight_slide.shapes.add_shape(1, 571499, 6100000, 11100000, 500000)
sum_bar.fill.solid()
sum_bar.fill.fore_color.rgb = SOL_BG
sum_bar.line.fill.background()

add_body_text(insight_slide, 671499, 6200000, 11000000, 300000,
              '💬 바이브코딩은 단순 도구를 넘어, 업무 사고방식 자체를 AI 협업 중심으로 전환하는 계기가 되었음',
              color=RED_ACCENT, size=171450, bold=True)

# Bottom bar
bottom_bar = insight_slide.shapes.add_shape(1, 0, H - 100000, W, 100000)
bottom_bar.fill.solid()
bottom_bar.fill.fore_color.rgb = RED
bottom_bar.line.fill.background()

# Move insight slide to before the 감사합니다 slide (last slide)
# Currently at end, want it just before the last original slide
move_slide(prs, len(prs.slides) - 1, len(prs.slides) - 2)
print(f'Insight slide moved to position {len(prs.slides)-2}')

print(f'\nFinal slide count: {len(prs.slides)}')
print('Saving...')
prs.save(DST)
print(f'Saved to {DST}')
